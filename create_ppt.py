from pptx import Presentation
from pptx.util import Inches
import os
import glob
import yaml
from PIL import Image
import sys

def load_config():
    config_file_path = os.path.join(os.path.dirname(__file__), "config.yaml")
    with open(config_file_path, "r", encoding="utf-8") as f:
        cfg_data = yaml.safe_load(f)
    return cfg_data

def compress_and_resize_image(input_path, output_path, target_width, target_height, quality=80):
    """
    调整图片大小到目标尺寸，并用指定的 JPEG 质量保存
    """
    with Image.open(input_path) as img:
        # 采用强制 resize 到目标尺寸，如有需要可保留图片比例
        img_resized = img.resize((target_width, target_height), Image.LANCZOS)
        img_resized.save(output_path, "JPEG", quality=quality)

def create_ppt(image_paths, ppt_config, compressed_folder):
    """
    根据图片列表及 PPT 配置生成幻灯片，每张幻灯片添加一张经过预处理的图片。
    返回生成的 PPT 文件的路径。
    """
    ppt_path = ppt_config.get("ppt_path", "output.pptx")
    slide_width_inches = ppt_config.get("slide_width", 16)
    slide_height_inches = ppt_config.get("slide_height", 9)
    
    # 初始化 PPTX 对象，并设置幻灯片大小
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    
    slide_layout_index = ppt_config.get("slide_layout_index", 6)
    blank_slide_layout = prs.slide_layouts[slide_layout_index]
    
    # 如果压缩图片目录不存在，则创建该目录
    if not os.path.exists(compressed_folder):
        os.makedirs(compressed_folder)
    
    # 计算目标图片尺寸（基于 96 DPI）
    dpi = 96
    target_width = int(slide_width_inches * dpi)
    target_height = int(slide_height_inches * dpi)
    
    compressed_image_paths = []
    for img_path in image_paths:
        base = os.path.basename(img_path)
        compressed_path = os.path.join(compressed_folder, base)
        compress_and_resize_image(img_path, compressed_path, target_width, target_height, quality=80)
        compressed_image_paths.append(compressed_path)
    
    # 使用压缩后的图片生成幻灯片，每张幻灯片添加一张图片
    for comp_img in compressed_image_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        slide.shapes.add_picture(comp_img, left, top, width=width, height=height)
    
    # 保存 PPT 文件到指定路径
    prs.save(ppt_path)
    print(f"PPT 已保存到 {ppt_path}")
    return ppt_path  # 返回生成的 PPT 文件路径

if __name__ == "__main__":
    # 判断是否传入了用户目录（例如 "1", "2", "3", …）作为参数
    if len(sys.argv) < 2:
        print("错误：请提供用户目录名称（例如 1、2、3 等）作为第一个参数。")
        sys.exit(1)
        
    user_dir_name = sys.argv[1]
    
    # 获取项目根目录，并构造 upload_history 内对应用户目录
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    upload_history_dir = os.path.join(BASE_DIR, "upload_history")
    user_dir = os.path.join(upload_history_dir, user_dir_name)
    
    # 定义图片来源目录为 user_dir 下的 results 文件夹
    results_folder = os.path.join(user_dir, "results")
    # 定义 PPT 输出目录为 user_dir 下的 generated_ppt 文件夹
    generated_ppt_folder = os.path.join(user_dir, "generated_ppt")
    if not os.path.exists(generated_ppt_folder):
        os.makedirs(generated_ppt_folder)
    
    # 设置压缩图片的临时目录，可放在 user_dir 下（例如：compressed_images）
    compressed_folder = os.path.join(user_dir, "compressed_images")
    
    # 重载配置文件（如果 config.yaml 中有 PPT 相关配置）
    cfg = load_config()
    ppt_config = cfg.get("PPT", {})
    # 固定 PPT 存储路径到 generated_ppt 目录下
    ppt_config["ppt_path"] = os.path.join(generated_ppt_folder, "output.pptx")
    
    # 收集 results 目录下所有 jpg/jpeg 图片，并排序
    image_paths = glob.glob(os.path.join(results_folder, "*.jpg")) + \
                  glob.glob(os.path.join(results_folder, "*.jpeg"))
    image_paths = sorted(image_paths)
    
    if not image_paths:
        print(f"错误：在 {results_folder} 下未找到处理后的图片。")
        sys.exit(1)
    
    # 生成 PPT，使用预处理好的图片
    ppt_path = create_ppt(image_paths, ppt_config, compressed_folder)
    print(f"生成的 PPT 路径: {ppt_path}")
